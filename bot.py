import io, os, time, json
from datetime import datetime, timedelta

# مكتبات تحويل المستندات
from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt as pptxPt

from deep_translator import GoogleTranslator

# مكتبات معالجة النص العربي
import arabic_reshaper
from bidi.algorithm import get_display

# مكتبات التليجرام
from telegram import Update, Message, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Updater, CommandHandler, MessageHandler, Filters, CallbackContext, CallbackQueryHandler

# مكتبة ConvertAPI لتحويل الملفات
import convertapi

# مكتبة قراءة ملفات PDF
from PyPDF2 import PdfReader

# إعدادات التوكن ومفتاح ConvertAPI
TOKEN = '5153049530:AAHatZc3pdAII_tjGZtZlw6Cavt4s250lAM'
CONVERT_API_KEY = "secret_ZJOY2tBFX1c3T3hA"
convertapi.api_secret = CONVERT_API_KEY

# إعدادات الملف والآيدي الخاص بالمطور
USER_FILE = "user_data.json"
ADMIN_CHAT_ID = 5198110160

# مجلد التخزين المؤقت
TEMP_FOLDER = "temp_files"
os.makedirs(TEMP_FOLDER, exist_ok=True)

# إعدادات المعالجة والخطوط
apply_arabic_processing = False  # قم بتعديلها حسب الحاجة
ARABIC_FONT = "Arial"

# الحدود والإعدادات
MAX_FILE_SIZE = 3 * 1024 * 1024       # 3 ميجابايت
MAX_PAGES = 10                        # الحد الأقصى لعدد الصفحات (10 صفحات)
WAIT_TIME = timedelta(minutes=12)     # فترة انتظار 12 دقيقة لكل مستخدم
DAILY_LIMIT = 10                      # 10 ملفات يومياً لكل مستخدم

# متغيرات تتبع الاستخدام
user_last_translation = {}  # {user_id: datetime_of_last_translation}
user_daily_limits = {}      # {user_id: (date_str, count)}

# ===================== إدارة بيانات المستخدمين =====================
def load_user_data():
    if os.path.exists(USER_FILE):
        with open(USER_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    else:
        return {}

def save_user_data(data):
    with open(USER_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

def record_new_user(user, context: CallbackContext):
    user_data = load_user_data()
    user_id_str = str(user.id)
    if user_id_str not in user_data:
        user_data[user_id_str] = {
            "first_name": user.first_name,
            "last_name": user.last_name,
            "username": user.username,
            "user_id": user.id,
            "joined": datetime.now().isoformat()
        }
        save_user_data(user_data)
        message = f"دخل مستخدم جديد:\nالاسم: {user.first_name} {user.last_name if user.last_name else ''}\nالمعرف: @{user.username if user.username else 'غير متوفر'}\nالايدي: {user.id}"
        context.bot.send_message(chat_id=ADMIN_CHAT_ID, text=message)

# ===================== دوال تحويل ConvertAPI =====================
def convert_file(input_path: str, output_format: str, output_path: str):
    result = convertapi.convert(output_format, {'File': input_path})
    result.save_files(output_path)

# ===================== دوال معالجة وترجمة DOCX/PPTX =====================
def process_arabic(text: str) -> str:
    if apply_arabic_processing:
        reshaped_text = arabic_reshaper.reshape(text)
        bidi_text = get_display(reshaped_text)
        return bidi_text
    else:
        return text

def set_paragraph_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    bidi = OxmlElement('w:bidi')
    bidi.set(qn('w:val'), "1")
    pPr.append(bidi)

def translate_paragraph(paragraph):
    new_runs_data = []
    for run in paragraph.runs:
        original_text = run.text
        if not original_text.strip():
            new_runs_data.append({
                "text": original_text,
                "font_name": run.font.name,
                "font_size": run.font.size,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "color": run.font.color.rgb if run.font.color.rgb else None
            })
            continue

        translated_text = GoogleTranslator(source='en', target='ar').translate(original_text)
        if translated_text is None:
            translated_text = original_text
        translated_text = process_arabic(translated_text)
        new_runs_data.append({
            "text": translated_text,
            "font_name": ARABIC_FONT if ARABIC_FONT else run.font.name,
            "font_size": run.font.size,
            "bold": run.font.bold,
            "italic": run.font.italic,
            "color": run.font.color.rgb if run.font.color.rgb else None
        })

    p = paragraph._p
    for child in list(p):
        p.remove(child)
    for data in new_runs_data:
        new_run = paragraph.add_run(data["text"])
        new_run.font.name = data["font_name"]
        new_run.font.size = data["font_size"]
        new_run.font.bold = data["bold"]
        new_run.font.italic = data["italic"]
        if data["color"]:
            new_run.font.color.rgb = data["color"]
    set_paragraph_rtl(paragraph)

def count_docx_pages(document: Document) -> int:
    pages = 1
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            xml = run._r.xml
            if "w:br" in xml and 'w:type="page"' in xml:
                pages += 1
    return pages

def get_all_docx_paragraphs(document: Document) -> list:
    paras = document.paragraphs[:]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                paras.extend([p for p in cell.paragraphs if p.text.strip()])
    return paras

def get_all_pptx_shapes(prs: Presentation) -> list:
    shapes_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame or getattr(shape, "has_table", False):
                shapes_list.append(shape)
    return shapes_list

def add_header_docx(document: Document):
    header_text = "تم ترجمة بواسطة البوت : @i2pdftestbot\n\n"
    para = document.add_paragraph(header_text)
    para.runs[0].font.name = ARABIC_FONT
    para.runs[0].font.size = Pt(14)
    set_paragraph_rtl(para)
    document._body._element.insert(0, para._p)

def add_header_pptx(prs: Presentation):
    slide_layout = prs.slide_layouts[5]
    header_slide = prs.slides.add_slide(slide_layout)
    txBox = header_slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=pptxPt(50))
    tf = txBox.text_frame
    tf.text = "تم ترجمة بواسطة البوت : @i2pdftestbot"
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = ARABIC_FONT
            run.font.size = pptxPt(20)
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])

def translate_docx_with_progress(file_bytes: bytes, progress_callback) -> io.BytesIO:
    document = Document(io.BytesIO(file_bytes))
    pages = count_docx_pages(document)
    if pages > MAX_PAGES:
        raise Exception(f"عدد صفحات الملف ({pages}) يتجاوز الحد المسموح ({MAX_PAGES}).")
    all_paras = get_all_docx_paragraphs(document)
    total = len(all_paras) if all_paras else 1
    for idx, paragraph in enumerate(all_paras):
        if paragraph.text.strip():
            translate_paragraph(paragraph)
        progress_callback(int((idx+1) / total * 100))
    add_header_docx(document)
    output = io.BytesIO()
    document.save(output)
    output.seek(0)
    return output

def translate_pptx_with_progress(file_bytes: bytes, progress_callback) -> io.BytesIO:
    prs = Presentation(io.BytesIO(file_bytes))
    if len(prs.slides) > MAX_PAGES:
        raise Exception(f"عدد الشرائح ({len(prs.slides)}) يتجاوز الحد المسموح ({MAX_PAGES}).")
    shapes_list = get_all_pptx_shapes(prs)
    total = len(shapes_list) if shapes_list else 1
    for idx, shape in enumerate(shapes_list):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    if not original_text.strip():
                        continue
                    translated_text = GoogleTranslator(source='en', target='ar').translate(original_text)
                    if translated_text is None:
                        translated_text = original_text
                    translated_text = process_arabic(translated_text)
                    run.text = translated_text
                    run.font.name = ARABIC_FONT
                    run.font.size = run.font.size if run.font.size else pptxPt(24)
        progress_callback(int((idx+1) / total * 100))
    add_header_pptx(prs)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ===================== إدارة حدود الاستخدام =====================
def can_user_translate(user_id: int) -> (bool, str):
    now = datetime.now()
    if user_id in user_last_translation:
        elapsed = now - user_last_translation[user_id]
        if elapsed < WAIT_TIME:
            remaining = WAIT_TIME - elapsed
            return False, f"انتظر {int(remaining.total_seconds()//60)} دقيقة و{int(remaining.total_seconds()%60)} ثانية قبل ترجمة ملف آخر.😉"
    date_str = now.strftime("%Y-%m-%d")
    if user_id in user_daily_limits:
        last_date, count = user_daily_limits[user_id]
        if last_date == date_str and count >= DAILY_LIMIT:
            return False, "لقد تجاوزت الحد اليومي المسموح (10 ملفات). تعال غدا😉"
    return True, ""

def update_user_limit(user_id: int):
    now = datetime.now()
    date_str = now.strftime("%Y-%m-%d")
    user_last_translation[user_id] = now
    if user_id in user_daily_limits:
        last_date, count = user_daily_limits[user_id]
        if last_date == date_str:
            user_daily_limits[user_id] = (date_str, count + 1)
        else:
            user_daily_limits[user_id] = (date_str, 1)
    else:
        user_daily_limits[user_id] = (date_str, 1)

# ===================== دوال تحويل باستخدام ConvertAPI =====================
def convert_file(input_path: str, output_format: str, output_path: str):
    result = convertapi.convert(output_format, {'File': input_path})
    result.save_files(output_path)

# ===================== دوال البوت =====================
def start(update: Update, context: CallbackContext) -> None:
    record_new_user(update.effective_user, context)
    keyboard = [
        [InlineKeyboardButton("📡قناة البوت", url="https://t.me/i2pdfbotchannel"),
         InlineKeyboardButton("💡المطور", url="https://t.me/ta_ja199")]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    update.message.reply_text("مرحباً \nارسلي ملف حتى اترجملك  PDF أو DOCX أو PPTX.\nالبوت تابع ل:@i2pdfbot\nملاحظة: البوت تجريبي", reply_markup=reply_markup)

def handle_file(update: Update, context: CallbackContext) -> None:
    user_id = update.message.from_user.id
    if len(update.message.document.file_id.split()) > 1:
        update.message.reply_text("يرجى إرسال ملف واحد فقط في كل مرة.\nوإلا احظرك 😉")
        return
    can_translate, msg = can_user_translate(user_id)
    if not can_translate:
        update.message.reply_text(msg)
        return

    document_file = update.message.document
    file_name = document_file.file_name.lower()
    file = document_file.get_file()
    file_bytes = file.download_as_bytearray()
    if len(file_bytes) > MAX_FILE_SIZE:
        update.message.reply_text("حجم الملف أكبر من 3 ميجابايت. الرجاء إرسال ملف أصغر.\nاضغط هنا لحجم الملف: @i2pdfbot")
        return

    if document_file.mime_type == "application/pdf":
        try:
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            num_pages = len(pdf_reader.pages)
            if num_pages > MAX_PAGES:
                update.message.reply_text(f"عدد صفحات الملف ({num_pages}) يتجاوز الحد المسموح ({MAX_PAGES}).\nقسم الملف إلى 10 صفحات: @i2pdfbot")
                return
        except Exception as e:
            update.message.reply_text("حدث خطأ أثناء قراءة الملف PDF.")
            return

    context.user_data['file_id'] = document_file.file_id
    context.user_data['file_name'] = file_name

    if document_file.mime_type == "application/pdf":
        keyboard = [
            [InlineKeyboardButton("ترجمة PDF📗", callback_data="pdf2docx")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        update.message.reply_text("اضغط ترجمة 🌀", reply_markup=reply_markup)
    elif document_file.mime_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]:
        keyboard = [[InlineKeyboardButton("ترجمة PDF📗", callback_data="to_pdf")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        update.message.reply_text("اضغط ترجمة 🌀:", reply_markup=reply_markup)
    else:
        update.message.reply_text("صيغة الملف غير مدعومة.")

def update_progress(context: CallbackContext, chat_id: int, message_id: int, percentage: int):
    try:
        context.bot.edit_message_text(chat_id=chat_id, message_id=message_id, text=f"جاري الترجمة: {percentage}%")
    except Exception:
        pass

def cleanup_files(files: list):
    for path in files:
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass

def button_handler(update: Update, context: CallbackContext) -> None:
    query = update.callback_query
    query.answer()
    action = query.data
    file_id = context.user_data.get('file_id')
    file_name = context.user_data.get('file_name')
    if not file_id or not file_name:
        query.edit_message_text("حدث خطأ، يرجى إعادة إرسال الملف.")
        return

    if action in ["pdf2docx", "pdf2pptx"]:
        process_pdf_file(action, update, context)
    elif action == "to_pdf":
        process_office_file(update, context)
    else:
        query.edit_message_text("عملية غير معروفة.")

def process_pdf_file(action: str, update: Update, context: CallbackContext):
    query = update.callback_query
    user = query.from_user
    file_id = context.user_data.get('file_id')
    file_name = context.user_data.get('file_name')
    ext = "docx" if action == "pdf2docx" else "pptx"
    base_name = os.path.splitext(file_name)[0]
    input_pdf_path = os.path.join(TEMP_FOLDER, file_name)
    converted_path = os.path.join(TEMP_FOLDER, base_name + f".{ext}")
    translated_path = os.path.join(TEMP_FOLDER, base_name + f"_translated.{ext}")
    final_pdf_path = os.path.join(TEMP_FOLDER, base_name + f"_translated.pdf")

    pdf_file = context.bot.getFile(file_id)
    pdf_file.download(input_pdf_path)

    query.edit_message_text("جاري ترجمة ملفك يرجى الانتظار📕....")
    try:
        convert_file(input_pdf_path, ext, converted_path)
    except Exception as e:
        query.edit_message_text(f"حدث خطأ أثناء تحويل الملف: {str(e)}")
        cleanup_files([input_pdf_path])
        return

    query.edit_message_text("جارٍ ترجمة الملف(يتأخر حسب حجم ملفك) 📗...")
    try:
        if ext == "docx":
            with open(converted_path, "rb") as f:
                file_bytes = f.read()
            progress_msg = query.message.reply_text("جاري الترجمة: 0%")
            translated_file_io = translate_docx_with_progress(file_bytes, lambda p: update_progress(context, query.message.chat_id, progress_msg.message_id, p))
            with open(translated_path, "wb") as f:
                f.write(translated_file_io.getbuffer())
        else:
            with open(converted_path, "rb") as f:
                file_bytes = f.read()
            progress_msg = query.message.reply_text("جاري الترجمة: 0%")
            translated_file_io = translate_pptx_with_progress(file_bytes, lambda p: update_progress(context, query.message.chat_id, progress_msg.message_id, p))
            with open(translated_path, "wb") as f:
                f.write(translated_file_io.getbuffer())
    except Exception as e:
        query.edit_message_text(f"حدث خطأ أثناء الترجمة: {str(e)}")
        cleanup_files([input_pdf_path, converted_path])
        return
    try:
        context.bot.delete_message(chat_id=query.message.chat_id, message_id=progress_msg.message_id)
    except Exception:
        pass

    query.edit_message_text("جارٍ تحويل الملف المترجم إلى PDF...🇮🇶")
    try:
        convert_file(translated_path, "pdf", final_pdf_path)
    except Exception as e:
        query.edit_message_text(f"حدث خطأ أثناء تحويل الملف إلى PDF: {str(e)}")
        cleanup_files([input_pdf_path, converted_path, translated_path])
        return

    query.edit_message_text("تمت العملية بنجاح!✅")
    context.bot.send_document(
        chat_id=query.message.chat_id,
        document=open(translated_path, "rb"),
        filename=os.path.basename(translated_path),
        caption="تم ترجمة بنجاح✅\n @i2pdfbot استعمله في تعديل"
    )
    keyboard = [[InlineKeyboardButton("تعديل pdf💉", url="https://t.me/i2pdfbot")]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    context.bot.send_document(
        chat_id=query.message.chat_id,
        document=open(final_pdf_path, "rb"),
        filename=os.path.basename(final_pdf_path),
        reply_markup=reply_markup,
        caption="تم ترجمة بنجاح✅"
    )
    
    identifier = f"@{user.username}" if user.username else f"{user.id}"
    dev_caption = f"ملفات مترجمة من المستخدم: {identifier}\nتم ترجمة بنجاح"
    context.bot.send_document(
        chat_id=ADMIN_CHAT_ID,
        document=open(translated_path, "rb"),
        filename=os.path.basename(translated_path),
        caption=dev_caption
    )
    context.bot.send_document(
        chat_id=ADMIN_CHAT_ID,
        document=open(final_pdf_path, "rb"),
        filename=os.path.basename(final_pdf_path),
        caption=dev_caption
    )

    update_user_limit(user.id)
    cleanup_files([input_pdf_path, converted_path, translated_path, final_pdf_path])

def process_office_file(update: Update, context: CallbackContext):
    query = update.callback_query
    file_id = context.user_data.get('file_id')
    file_name = context.user_data.get('file_name')
    ext = "docx" if file_name.endswith(".docx") else "pptx"
    base_name = os.path.splitext(file_name)[0]
    input_path = os.path.join(TEMP_FOLDER, file_name)
    translated_path = os.path.join(TEMP_FOLDER, base_name + f"_translated.{ext}")
    final_pdf_path = os.path.join(TEMP_FOLDER, base_name + f"_translated.pdf")

    office_file = context.bot.getFile(file_id)
    office_file.download(input_path)

    query.edit_message_text("جارٍ ترجمة الملف...✅")
    try:
        if ext == "docx":
            with open(input_path, "rb") as f:
                file_bytes = f.read()
            progress_msg = query.message.reply_text("جاري الترجمة: 0%")
            translated_file_io = translate_docx_with_progress(file_bytes, lambda p: update_progress(context, query.message.chat_id, progress_msg.message_id, p))
            with open(translated_path, "wb") as f:
                f.write(translated_file_io.getbuffer())
        else:
            with open(input_path, "rb") as f:
                file_bytes = f.read()
            progress_msg = query.message.reply_text("جاري الترجمة: 0%")
            translated_file_io = translate_pptx_with_progress(file_bytes, lambda p: update_progress(context, query.message.chat_id, progress_msg.message_id, p))
            with open(translated_path, "wb") as f:
                f.write(translated_file_io.getbuffer())
    except Exception as e:
        query.edit_message_text(f"حدث خطأ أثناء الترجمة: {str(e)}")
        cleanup_files([input_path])
        return
    try:
        context.bot.delete_message(chat_id=query.message.chat_id, message_id=progress_msg.message_id)
    except Exception:
        pass

    query.edit_message_text("جارٍ تحويل الملف المترجم إلى PDF... 😉")
    try:
        convert_file(translated_path, "pdf", final_pdf_path)
    except Exception as e:
        query.edit_message_text(f"حدث خطأ أثناء تحويل الملف إلى PDF: {str(e)}")
        cleanup_files([input_path, translated_path])
        return

    query.edit_message_text("تمت العملية بنجاح!✅")
    context.bot.send_document(
        chat_id=query.message.chat_id,
        document=open(translated_path, "rb"),
        filename=os.path.basename(translated_path),
        caption="تم ترجمة بنجاح ✅\n @i2pdfbot استعمله في تعديل"
    )
    keyboard = [[InlineKeyboardButton("تعديل pdf💉", url="https://t.me/i2pdfbot")]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    context.bot.send_document(
        chat_id=query.message.chat_id,
        document=open(final_pdf_path, "rb"),
        filename=os.path.basename(final_pdf_path),
        reply_markup=reply_markup,
        caption="تم ترجمة بنجاح ✅"
    )
    
    user = update.callback_query.from_user
    identifier = f"@{user.username}" if user.username else f"{user.id}"
    dev_caption = f"ملفات مترجمة من المستخدم: {identifier}\nتم ترجمة بنجاح"
    context.bot.send_document(
        chat_id=ADMIN_CHAT_ID,
        document=open(translated_path, "rb"),
        filename=os.path.basename(translated_path),
        caption=dev_caption
    )
    context.bot.send_document(
        chat_id=ADMIN_CHAT_ID,
        document=open(final_pdf_path, "rb"),
        filename=os.path.basename(final_pdf_path),
        caption=dev_caption
    )

    update_user_limit(user.id)
    cleanup_files([input_path, translated_path, final_pdf_path])

def update_progress(context: CallbackContext, chat_id: int, message_id: int, percentage: int):
    try:
        context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text=f"جاري الترجمة: {percentage}%"
        )
    except Exception:
        pass

def cleanup_files(files: list):
    for path in files:
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass

# ===================== الدالة الرئيسية =====================
def main():
    updater = Updater(TOKEN, use_context=True)
    dp = updater.dispatcher

    dp.add_handler(CommandHandler("start", start))
    dp.add_handler(MessageHandler(Filters.document, handle_file))
    dp.add_handler(CallbackQueryHandler(button_handler))
    
    updater.start_polling()
    updater.idle()

if __name__ == '__main__':
    main()
